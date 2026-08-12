#!/usr/bin/env python3
"""Turn a built deck into a 16:9 PowerPoint file, one rendered slide per page.

The decks are HTML, and their layout leans on things PowerPoint has no notion of
(CSS grid, webfont metrics, a real product screenshot bled onto the ground). So
this does not translate the deck into PowerPoint shapes - it renders each slide
in a headless browser at 2x and lays the images into a 13.333x7.5in deck, which
is what a designed deck survives as. Text is therefore not editable in the
result; edit src/ and rebuild.

Speaker notes ARE carried across: each slide's data-t/data-n become the notes
page, so presenter view works.

Extra requirements, deliberately not project dependencies:
    pip install playwright python-pptx && playwright install chromium

Usage:
    python tools/embed_assets.py src/deck-pitch.html      # build dist/ first
    python tools/build_pptx.py dist/deck-pitch.html
"""

from __future__ import annotations

import re
import sys
import tempfile
from pathlib import Path

SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5


def shoot(deck: Path, out_dir: Path) -> list[Path]:
    from playwright.sync_api import sync_playwright

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(
            viewport={"width": 1280, "height": 720}, device_scale_factor=2
        )
        page.goto(deck.as_uri() + "#1")
        page.wait_for_function("document.fonts.ready.then(()=>true)")
        # undo the fit-to-window scale so the reel renders at native size, and
        # hide the presenter chrome - it is fixed to the viewport, so an element
        # screenshot of the reel would otherwise capture it lying on top
        page.evaluate("document.getElementById('reel').style.transform='none'")
        page.add_style_tag(content="#bar,#hud{display:none!important}")
        page.wait_for_timeout(400)
        count = page.evaluate("document.querySelectorAll('.slide').length")
        reel = page.locator("#reel")
        shots = []
        for i in range(1, count + 1):
            page.wait_for_timeout(220)
            shot = out_dir / f"slide-{i:02d}.png"
            reel.screenshot(path=str(shot))
            shots.append(shot)
            page.keyboard.press("ArrowRight")
        browser.close()
    return shots


def notes_for(src: Path) -> list[tuple[str, str]]:
    """(title, speaker note) per slide, read from the deck source."""
    if not src.exists():
        return []
    html = src.read_text(encoding="utf-8")
    out = []
    for chunk in re.findall(r'<section class="slide[^"]*"(.*?)>', html, re.S):
        vals = []
        for name in ("data-t", "data-n"):
            m = re.search(rf'{name}="(.*?)"', chunk, re.S)
            text = re.sub(r"\s+", " ", m.group(1)).strip() if m else ""
            for a, b in (("&amp;", "&"), ("&mdash;", "—"), ("&nbsp;", " "), ("&quot;", '"')):
                text = text.replace(a, b)
            vals.append(text)
        out.append((vals[0], vals[1]))
    return out


def main(argv: list[str]) -> int:
    if not argv:
        print(__doc__)
        return 2

    from pptx import Presentation
    from pptx.util import Inches

    deck = Path(argv[0]).resolve()
    if not deck.exists():
        print(f"!! {deck} not found - run tools/embed_assets.py first")
        return 1
    src = deck.parents[1] / "src" / deck.name
    notes = notes_for(src)

    with tempfile.TemporaryDirectory() as tmp:
        shots = shoot(deck, Path(tmp))

        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(SLIDE_W_IN), Inches(SLIDE_H_IN)
        blank = prs.slide_layouts[6]
        for i, shot in enumerate(shots):
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                str(shot), 0, 0, width=Inches(SLIDE_W_IN), height=Inches(SLIDE_H_IN)
            )
            if i < len(notes):
                title, note = notes[i]
                body = f"{title}\n\n{note}".strip() if note else title
                if body:
                    slide.notes_slide.notes_text_frame.text = body

        out = deck.with_suffix(".pptx")
        prs.save(out)
    print(f"ok {len(shots)} slides -> {out.parent.name}/{out.name} "
          f"({out.stat().st_size/1024/1024:.1f} MB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
